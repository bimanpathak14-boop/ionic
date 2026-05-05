"""PowerPoint Presentation Handler"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PowerPointHandler:
    def __init__(self):
        self.output_dir = os.path.expanduser('~/Documents/PocketAI')
        os.makedirs(self.output_dir, exist_ok=True)

    def create(self, title='Untitled Presentation', slides=None,
               theme='modern', author='Pocket AI Office', **kwargs):
        """Create a PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        colors = self._get_theme_colors(theme)

        # Title slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide, colors['bg'])

        # Title text
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = colors['title']
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub = tf.add_paragraph()
        sub.text = f'Created by {author}'
        sub.font.size = Pt(18)
        sub.font.color.rgb = colors['subtitle']
        sub.alignment = PP_ALIGN.CENTER

        # Content slides
        if slides:
            for slide_data in slides:
                self._add_content_slide(prs, slide_data, colors)

        # Save with unique filename if busy
        base_filename = self._safe_filename(title)
        filename = base_filename + '.pptx'
        filepath = os.path.join(self.output_dir, filename)
        
        counter = 1
        while os.path.exists(filepath):
            try:
                prs.save(filepath)
                break
            except PermissionError:
                filename = f"{base_filename}_{counter}.pptx"
                filepath = os.path.join(self.output_dir, filename)
                counter += 1
            except Exception:
                import datetime
                filename = f"{base_filename}_{datetime.datetime.now().strftime('%H%M%S')}.pptx"
                filepath = os.path.join(self.output_dir, filename)
                prs.save(filepath)
                break
        else:
            prs.save(filepath)
        
        # Open the file automatically
        try:
            os.startfile(filepath)
        except:
            pass

        return {
            'data': {'message': f'Presentation "{title}" created with {len(slides or []) + 1} slides',
                     'path': filepath},
            'filesCreated': [{'name': filename, 'path': filepath, 'type': 'presentation',
                              'mimeType': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'size': os.path.getsize(filepath)}],
        }

    def _add_content_slide(self, prs, slide_data, colors):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_bg(slide, colors['bg'])

        heading = slide_data.get('heading', '')
        subtitle = slide_data.get('subtitle', '')
        bullets = slide_data.get('bullets', [])
        content = slide_data.get('content', '')

        # Heading
        if heading:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = heading
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = colors['title']

        # Subtitle (new)
        if subtitle:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(22)
            p.font.color.rgb = colors['subtitle']

        # Bullet points
        if bullets:
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f' {bullet}'
                p.font.size = Pt(20)
                p.font.color.rgb = colors['text']
                p.space_after = Pt(12)
                p.level = 0

        # Plain content
        elif content:
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(18)
            p.font.color.rgb = colors['text']

    def _set_slide_bg(self, slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _get_theme_colors(self, theme):
        themes = {
            'modern': {
                'bg': RGBColor(0x1a, 0x1a, 0x2e),
                'title': RGBColor(0xff, 0xff, 0xff),
                'subtitle': RGBColor(0x99, 0x99, 0xbb),
                'text': RGBColor(0xdd, 0xdd, 0xee),
            },
            'corporate': {
                'bg': RGBColor(0xff, 0xff, 0xff),
                'title': RGBColor(0x1a, 0x3c, 0x6e),
                'subtitle': RGBColor(0x66, 0x66, 0x88),
                'text': RGBColor(0x33, 0x33, 0x44),
            },
            'creative': {
                'bg': RGBColor(0x0d, 0x0d, 0x1a),
                'title': RGBColor(0x66, 0x7e, 0xea),
                'subtitle': RGBColor(0x76, 0x4b, 0xa2),
                'text': RGBColor(0xcc, 0xcc, 0xee),
            },
        }
        return themes.get(theme, themes['modern'])

    def _safe_filename(self, name):
        return "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).strip()
